#!/usr/bin/env python3
"""Build an English WNRT portfolio PowerPoint (8–12 slides).

Usage:
  python scripts/build_wnrt_portfolio_pptx.py
  python scripts/build_wnrt_portfolio_pptx.py --out docs/wnrt-portfolio-deck.pptx
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Slate + teal (avoid default purple / cream-serif marketing look)
NAVY = RGBColor(0x0F, 0x1C, 0x2E)
SLATE = RGBColor(0x1E, 0x2A, 0x3A)
TEAL = RGBColor(0x1A, 0x9B, 0x8E)
TEAL_LIGHT = RGBColor(0xE6, 0xF5, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6B, 0x7C, 0x8F)
BODY = RGBColor(0x2C, 0x3A, 0x4B)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xF9)

REPO = Path(__file__).resolve().parents[1]
DEFAULT_OUT = REPO / "docs" / "wnrt-portfolio-deck.pptx"


def _set_run(run, text: str, *, size: int, bold: bool = False, color: RGBColor = BODY) -> None:
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _fill_shape(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _banner(slide, title: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    _fill_shape(bar, NAVY)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(1.15), Inches(13.333), Inches(0.08))
    _fill_shape(accent, TEAL)

    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12), Inches(0.55))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    _set_run(run, title, size=26, bold=True, color=WHITE)
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.55), Inches(0.72), Inches(12), Inches(0.35))
        p2 = box2.text_frame.paragraphs[0]
        r2 = p2.add_run()
        _set_run(r2, subtitle, size=12, color=TEAL_LIGHT)


def _footer(slide, page: int, total: int) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.15), Inches(10), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    _set_run(run, "WNRT · Technology Risk & Control Analytics Platform  ·  Portfolio deck", size=10, color=MUTED)
    box2 = slide.shapes.add_textbox(Inches(11.5), Inches(7.15), Inches(1.3), Inches(0.3))
    p2 = box2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    _set_run(r2, f"{page} / {total}", size=10, color=MUTED)


def _bullets(slide, left, top, width, height, items: list[str], *, size: int = 16) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(8 if i else 0)
        run = p.add_run()
        _set_run(run, f"•  {item}", size=size, color=BODY)


def _card(slide, left, top, width, height, title: str, body: str) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill_shape(shape, LIGHT_BG)
    shape.adjustments[0] = 0.08
    tbox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.35), Inches(0.35))
    p = tbox.text_frame.paragraphs[0]
    run = p.add_run()
    _set_run(run, title, size=14, bold=True, color=TEAL)
    bbox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.5), width - Inches(0.35), height - Inches(0.65))
    tf = bbox.text_frame
    tf.word_wrap = True
    p2 = tf.paragraphs[0]
    r2 = p2.add_run()
    _set_run(r2, body, size=13, color=BODY)


def build(out: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 10

    # 1 Title
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    _fill_shape(bg, NAVY)
    accent = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    _fill_shape(accent, TEAL)
    box = s.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    _set_run(run, "Windows Network Recovery Toolkit", size=36, bold=True, color=WHITE)
    box2 = s.shapes.add_textbox(Inches(0.9), Inches(3.0), Inches(11.5), Inches(0.6))
    p2 = box2.text_frame.paragraphs[0]
    r2 = p2.add_run()
    _set_run(r2, "Technology Risk & Control Analytics Platform", size=22, color=TEAL)
    box3 = s.shapes.add_textbox(Inches(0.9), Inches(3.8), Inches(11), Inches(1.2))
    tf = box3.text_frame
    tf.word_wrap = True
    p3 = tf.paragraphs[0]
    r3 = p3.add_run()
    _set_run(
        r3,
        "Evidence → classification → control tests → policy-gated remediation preview → audit custody.\n"
        "Portfolio deck  ·  English  ·  Not antivirus, EDR, or autonomous repair software.",
        size=15,
        color=TEAL_LIGHT,
    )
    _footer(s, 1, total)

    # 2 Problem
    s = prs.slides.add_slide(blank)
    _banner(s, "The problem", "Endpoints can look “online” while browsers and SSO fail")
    _bullets(
        s,
        Inches(0.7),
        Inches(1.6),
        Inches(11.8),
        Inches(4.8),
        [
            "Dead localhost WinINET proxies (e.g. 127.0.0.1:59081) with no listener",
            "WinINET vs WinHTTP drift — stacks disagree; ping/DNS still succeed",
            "TLS path mismatches and intermittent rewrites without registry-writer proof",
            "Operators guess Wi‑Fi / VPN / “virus” without reproducible evidence",
            "Risky “fix scripts” skip governance: no preview, no audit chain, no limitations",
        ],
        size=18,
    )
    _footer(s, 2, total)

    # 3 What it is / isn't
    s = prs.slides.add_slide(blank)
    _banner(s, "What WNRT is — and is not", "Positioning for FAANG / Big 4 / risk reviewers")
    _card(
        s,
        Inches(0.55),
        Inches(1.6),
        Inches(5.9),
        Inches(4.6),
        "We build",
        "Deterministic Windows endpoint evidence collection.\n\n"
        "Incident classification with proof tiers (T0–T5) and limitations[].\n\n"
        "Control testing, policy gates, dry-run remediation previews.\n\n"
        "Hash-chained audit custody and governance / Power BI exports.",
    )
    _card(
        s,
        Inches(6.8),
        Inches(1.6),
        Inches(5.9),
        Inches(4.6),
        "We do not claim",
        "Antivirus, EDR, XDR, or malware attribution.\n\n"
        "Autonomous remediation or AI-authorized apply.\n\n"
        "Formal audit opinions — reports are management information.\n\n"
        "Default process kill, firewall reset, or adapter disable.",
    )
    _footer(s, 3, total)

    # 4 Pipeline
    s = prs.slides.add_slide(blank)
    _banner(s, "Governance pipeline", "Stages stay separate — observation is not proof")
    stages = [
        ("1 Observe", "Registry, listeners,\nprobes, fixtures"),
        ("2 Hypothesize", "Triage labels +\nordinal confidence"),
        ("3 Prove", "Evidence tiers\nT0–T5"),
        ("4 Policy", "Allow / deny +\nconfirmation tokens"),
        ("5 Remediate", "Preview-first;\ndry-run default"),
        ("6 Audit", "JSONL custody +\ntip verify"),
    ]
    x = 0.45
    for title, body in stages:
        shape = s.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(1.95), Inches(3.2)
        )
        _fill_shape(shape, LIGHT_BG)
        shape.adjustments[0] = 0.1
        t = s.shapes.add_textbox(Inches(x + 0.1), Inches(2.4), Inches(1.75), Inches(0.7))
        p = t.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        _set_run(run, title, size=14, bold=True, color=TEAL)
        b = s.shapes.add_textbox(Inches(x + 0.1), Inches(3.3), Inches(1.75), Inches(1.8))
        tf = b.text_frame
        tf.word_wrap = True
        p2 = tf.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        _set_run(r2, body, size=12, color=BODY)
        x += 2.1
    note = s.shapes.add_textbox(Inches(0.55), Inches(5.7), Inches(12), Inches(0.8))
    p = note.text_frame.paragraphs[0]
    run = p.add_run()
    _set_run(
        run,
        "Correlation ≠ causation · Classification ≠ accusation · Policy allow ≠ safety guarantee",
        size=14,
        bold=True,
        color=SLATE,
    )
    _footer(s, 4, total)

    # 5 Safety
    s = prs.slides.add_slide(blank)
    _banner(s, "Safety model", "Remediation is gated; humans authorize apply")
    _bullets(
        s,
        Inches(0.7),
        Inches(1.55),
        Inches(11.8),
        Inches(5.0),
        [
            "Dry-run / preview by default (`--dry-run true`)",
            "Typed confirmation tokens for risky WinINET changes (e.g. DISABLE_WININET_PROXY)",
            "Blocked by default: process kill, firewall reset, adapter disable",
            "CI safety contracts: policy, classifier wording, dry-run API defaults",
            "Prefer deterministic fixtures over live-host speculation",
            "Preserve limitations[] in every operator-facing output",
        ],
        size=17,
    )
    _footer(s, 5, total)

    # 6 Architecture
    s = prs.slides.add_slide(blank)
    _banner(s, "Architecture snapshot", "Repo map reviewers can navigate in minutes")
    rows = [
        ("windows_network_toolkit/", "Primary CLI — proxy-status, diagnose, governance-report, agent"),
        ("src/platform_core/", "Policy, evidence tiers, remediation preview, audit writers"),
        ("src/proxy_drift/", "Startup observability, guardian, DNS health, evidence bundle"),
        ("backend/", "FastAPI — /trisk/*, /platform/*, enterprise routes"),
        ("tests/fixtures/", "Deterministic cases (e.g. dead proxy 59081)"),
        ("analytics/powerbi/", "Star-schema exports and interview-ready blueprints"),
    ]
    y = 1.55
    for left, right in rows:
        shape = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(y), Inches(12.2), Inches(0.72))
        _fill_shape(shape, LIGHT_BG)
        shape.adjustments[0] = 0.15
        t1 = s.shapes.add_textbox(Inches(0.75), Inches(y + 0.18), Inches(3.6), Inches(0.4))
        p = t1.text_frame.paragraphs[0]
        run = p.add_run()
        _set_run(run, left, size=14, bold=True, color=TEAL)
        t2 = s.shapes.add_textbox(Inches(4.4), Inches(y + 0.18), Inches(8.0), Inches(0.4))
        p2 = t2.text_frame.paragraphs[0]
        r2 = p2.add_run()
        _set_run(r2, right, size=14, color=BODY)
        y += 0.82
    _footer(s, 6, total)

    # 7 Case study
    s = prs.slides.add_slide(blank)
    _banner(s, "Case study: dead localhost proxy", "CASE_1_DEAD_WININET_PROXY · fixture-backed")
    _card(
        s,
        Inches(0.55),
        Inches(1.55),
        Inches(4.0),
        Inches(4.7),
        "Symptom",
        "ERR_PROXY_CONNECTION_FAILED / SSO timeouts.\n\n"
        "Ping and DNS succeed — host looks online.\n\n"
        "Users blame Wi‑Fi, VPN, or malware without evidence.",
    )
    _card(
        s,
        Inches(4.75),
        Inches(1.55),
        Inches(4.0),
        Inches(4.7),
        "Evidence (read-only)",
        "WinINET: ProxyEnable=1 → 127.0.0.1:59081\n\n"
        "WinHTTP: direct\n\n"
        "No listener on 59081\n\n"
        "Direct HTTPS OK; proxy path fails\n\n"
        "Proof tier T2 — not writer proof (T5)",
    )
    _card(
        s,
        Inches(8.95),
        Inches(1.55),
        Inches(3.8),
        Inches(4.7),
        "Outcome",
        "Class: DEAD_PROXY_CONFIG\n\n"
        "Reliability risk — not a malware verdict\n\n"
        "Remediation: preview + typed confirm\n\n"
        "Audit JSONL + governance export",
    )
    _footer(s, 7, total)

    # 8 Audience paths
    s = prs.slides.add_slide(blank)
    _banner(s, "Reviewer paths", "Same platform — different entry docs")
    _card(
        s,
        Inches(0.5),
        Inches(1.55),
        Inches(4.0),
        Inches(4.7),
        "FAANG / SRE / Platform",
        "docs/faang-platform-review.md\n\n"
        "State machine · API examples\n\n"
        "Fleet simulate & Docker demo\n\n"
        "Emphasis: contracts, stages, scale story",
    )
    _card(
        s,
        Inches(4.7),
        Inches(1.55),
        Inches(4.0),
        Inches(4.7),
        "Big 4 / Tech risk / Audit",
        "docs/big4-interview-defense.md\n\n"
        "Control matrix · sample governance report\n\n"
        "Evidence tiers & limitations[]\n\n"
        "Emphasis: defendable wording",
    )
    _card(
        s,
        Inches(8.9),
        Inches(1.55),
        Inches(3.9),
        Inches(4.7),
        "Power BI / Analytics",
        "docs/powerbi-interview-story.md\n\n"
        "Star-schema CSV exports\n\n"
        "Committee-ready KPIs\n\n"
        "Emphasis: explainable metrics",
    )
    _footer(s, 8, total)

    # 9 Demo proof
    s = prs.slides.add_slide(blank)
    _banner(s, "3-minute proof points", "Commands stay read-only unless you confirm")
    _bullets(
        s,
        Inches(0.7),
        Inches(1.55),
        Inches(11.8),
        Inches(5.2),
        [
            "python -m windows_network_toolkit proxy-status --fixture dead_proxy_60505.json",
            "python -m windows_network_toolkit proxy-disable --dry-run true",
            "python -m windows_network_toolkit reviewer-demo --mode mixed",
            "python -m windows_network_toolkit fleet-simulate --scenario mixed_proxy_failures --endpoints 100 --seed 42",
            "python -m windows_network_toolkit audit verify .audit/canonical_custody.jsonl --check-tip",
            "pytest -q tests/test_policy_safety_contract.py",
            "Docs: interview-demo-3min.md · PORTFOLIO.md · one-page-case-study-dead-proxy.md",
        ],
        size=16,
    )
    _footer(s, 9, total)

    # 10 Close
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    _fill_shape(bg, NAVY)
    accent = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    _fill_shape(accent, TEAL)
    box = s.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11.5), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    _set_run(run, "Takeaway", size=18, bold=True, color=TEAL)
    box2 = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.2))
    tf = box2.text_frame
    tf.word_wrap = True
    p2 = tf.paragraphs[0]
    r2 = p2.add_run()
    _set_run(
        r2,
        "WNRT turns messy Windows endpoint reliability signals into explainable classifications,\n"
        "control results, policy-gated previews, and audit-backed governance exports —\n"
        "without pretending to be antivirus or an autonomous fixer.",
        size=18,
        color=WHITE,
    )
    box3 = s.shapes.add_textbox(Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.2))
    tf3 = box3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    r3 = p3.add_run()
    _set_run(
        r3,
        "Next: live fixture demo  ·  governance report walkthrough  ·  fleet-simulate\n"
        "Start: README.md  ·  PORTFOLIO.md  ·  docs/ONBOARDING.md",
        size=14,
        color=TEAL_LIGHT,
    )
    _footer(s, 10, total)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def main() -> int:
    parser = argparse.ArgumentParser(description="Build WNRT English portfolio PPTX")
    parser.add_argument("--out", type=Path, default=DEFAULT_OUT)
    args = parser.parse_args()
    path = build(args.out.resolve())
    print(f"Wrote {path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
